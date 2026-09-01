import os
import pandas as pd
import logging
from copy import deepcopy
from datetime import datetime
import shutil
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def clone_slide(prs, slide):
    """
    Duplicate a slide.
    """

    blank_layout = prs.slide_layouts[6]

    new_slide = prs.slides.add_slide(
        blank_layout
    )

    for shape in slide.shapes:
        el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(
            el,
            "p:extLst"
        )

    return new_slide

def initialize_ppt_run():
    """
    Create dedicated PPT run folder.
    """

    base_dir = os.getcwd()

    timestamp = datetime.now().strftime(
        "%Y%m%d_%H%M%S_%f"
    )

    run_dir = os.path.join(
        base_dir,
        "data",
        "runs",
        f"ppt_run_{timestamp}"
    )

    output_dir = os.path.join(
        run_dir,
        "outputs"
    )

    os.makedirs(
        output_dir,
        exist_ok=True
    )

    return run_dir, output_dir

def update_cohort_text(
    slide,
    cohort_id
):

    for shape in slide.shapes:

        if (
            hasattr(shape, "name")
            and shape.name.lower()
            == "txt_cohort_id"
        ):

            if not shape.has_text_frame:
                return

            tf = shape.text_frame

            if (
                tf.paragraphs
                and tf.paragraphs[0].runs
            ):

                tf.paragraphs[0].runs[0].text = (
                    cohort_id
                )

            else:

                tf.text = cohort_id

            return

def load_ppt_rows(
    run_id
):

    processing_driver_file = os.path.join(
        os.getcwd(),
        "data",
        "input",
        "params",
        "processing_driver.csv"
    )

    df = pd.read_csv(
        processing_driver_file
    )
    
    return df[
        pd.to_numeric(
            df["enabled"],
            errors="coerce"
        ).fillna(0).astype(int) == 1
    ]

def add_template_slide(
    master_prs,
    template_slide
):

    blank_layout = master_prs.slide_layouts[6]

    slide = master_prs.slides.add_slide(
        blank_layout
    )

    for shape in template_slide.shapes:

        el = deepcopy(shape.element)

        slide.shapes._spTree.insert_element_before(
            el,
            "p:extLst"
        )

    return slide

def build_powerpoint(
    run_id,
    output_file=None
):
    rows = load_ppt_rows(
        run_id
    )

    logging.info(
        f"[ppt] Loaded {len(rows)} enabled rows"
    )

    logging.info(
        f"[ppt] Cohorts found: "
        f"{rows['cohort_id'].nunique()}"
    )

    ppt_run_dir, ppt_output_dir = (
        initialize_ppt_run()
    )

    logging.info(
        f"[ppt] PPT run initialized: "
        f"{ppt_run_dir}"
    )

    master_prs = None

    for cohort_id, cohort_rows in rows.groupby(
        "cohort_id"
    ):

        #
        # backward compatibility
        #
        if "ppt_slide_group" not in cohort_rows.columns:

            cohort_rows["ppt_slide_group"] = (
                cohort_rows["visual_id"]
            )

        if "ppt_slide_group_order" not in cohort_rows.columns:

            cohort_rows["ppt_slide_group_order"] = (
                cohort_rows["ppt_visual_order"]
            )

        #
        # sort slide groups
        #
        cohort_rows = cohort_rows.sort_values(
            [
                "ppt_slide_group_order",
                "ppt_visual_order"
            ]
        )

        #
        # one slide per group
        #
        for slide_group, slide_rows in (
            cohort_rows.groupby(
                "ppt_slide_group"
            )
        ):

            slide_rows = slide_rows.sort_values(
                "ppt_visual_order"
            )

            first_row = slide_rows.iloc[0]

            template_file = os.path.join(
                os.getcwd(),
                "data",
                "input",
                "templates",
                first_row["ppt_template"]
            )

            #
            # validate templates
            #
            template_values = (
                slide_rows["ppt_template"]
                .dropna()
                .astype(str)
                .unique()
            )

            if len(template_values) != 1:

                logging.error(
                    "[ppt] Slide group "
                    f"{slide_group} "
                    "contains multiple templates"
                )

                continue

            if master_prs is None:

                master_prs = Presentation(
                    template_file
                )

            template_prs = Presentation(
                template_file
            )

            template_slide = (
                template_prs.slides[0]
            )

            base_dir = os.getcwd()

            outputs_dir = os.path.join(
                base_dir,
                "data",
                "runs",
                run_id,
                "outputs"
            )

            if not os.path.exists(outputs_dir):

                raise FileNotFoundError(
                    outputs_dir
                )

            cohort_dir = os.path.join(
                outputs_dir,
                cohort_id
            )

            if not os.path.exists(cohort_dir):

                logging.warning(
                    f"[ppt] Skipping cohort {cohort_id}. "
                    f"Output folder not found: {cohort_dir}"
                )

                continue

            #
            # CREATE SINGLE SLIDE
            #
            slide = add_template_slide(
                master_prs,
                template_slide
            )

            update_cohort_text(
                slide,
                cohort_id
            )

            #
            # APPLY ALL VISUALS
            #
            for _, row in (
                slide_rows.iterrows()
            ):

                visual_id = row[
                    "visual_id"
                ]

                params = row.to_dict()

                populate_slide(
                    slide,
                    cohort_dir,
                    params
                )

                logging.info(
                    "[ppt] Applied "
                    f"{visual_id} "
                    f"to group "
                    f"{slide_group}"
                )

            remove_unwanted_placeholders(
                slide
            )

            logging.info(
                "[ppt] Added slide "
                f"group={slide_group} "
                f"cohort={cohort_id}"
            )

    if master_prs is None:

        raise RuntimeError(
            "[ppt] No slides were generated. "
            "Verify processing_driver.csv, "
            "templates, and cohort output folders."
        )

    if len(master_prs.slides) > 1:

        r_id = master_prs.slides._sldIdLst[0].rId

        master_prs.part.drop_rel(
            r_id
        )

        del master_prs.slides._sldIdLst[0]

    if output_file is None:

        output_file = os.path.join(
            ppt_output_dir,
            "powerpoint_poc.pptx"
        )

    master_prs.save(
        output_file
    )

    logging.info(
        f"[ppt] Generated "
        f"{len(master_prs.slides)} slides"
    )

    logging.info(
        f"[ppt] Saved: {output_file}"
    )

def get_shape_by_name(
    slide,
    shape_name
):

    for shape in slide.shapes:

        if (
            hasattr(shape, "name")
            and shape.name.lower()
            == shape_name.lower()
        ):
            return shape

    return None

def get_assets(
    cohort_dir,
    params
):

    assets = {}

    asset_numbers = set()

    for key in params.keys():

        if (
            str(key).startswith("ppt_asset_")
            and str(key).endswith("_keyword")
        ):

            parts = key.split("_")

            if len(parts) >= 3:
                asset_numbers.add(parts[2])

    for num in sorted(asset_numbers):

        keyword = params.get(
            f"ppt_asset_{num}_keyword"
        )

        placeholder = params.get(
            f"ppt_asset_{num}_placeholder"
        )

        if pd.isna(keyword) or pd.isna(placeholder):
            continue

        for filename in os.listdir(cohort_dir):

            lower = filename.lower()

            if not lower.endswith(".png"):
                continue

            keyword = str(keyword).strip().lower()

            if keyword in lower:

                assets[placeholder] = os.path.join(
                    cohort_dir,
                    filename
                )

                break

    return assets


def replace_picture(
    slide,
    shape_name,
    image_path
):

    if not image_path:
        logging.warning(
            f"[ppt] Missing image for "
            f"{shape_name}"
        )
        return

    shape = get_shape_by_name(
        slide,
        shape_name
    )

    if shape is None:

        logging.warning(
            f"[ppt] Placeholder not found: "
            f"{shape_name}"
        )

        return

    left = shape.left
    top = shape.top
    width = shape.width
    height = shape.height

    element = shape._element
    element.getparent().remove(
        element
    )

    new_pic = slide.shapes.add_picture(
        image_path,
        left,
        top,
        width=width,
        height=height
    )

    new_pic.name = shape_name

def populate_slide(
    slide,
    cohort_dir,
    params
):

    assets = get_assets(
        cohort_dir,
        params
    )

    for placeholder_name, image_path in assets.items():

        replace_picture(
            slide,
            placeholder_name,
            image_path
        )

def remove_unwanted_placeholders(slide):

    remove_prefixes = [
        "title",
        "text placeholder"
    ]

    remove_names = [
        "picture 58",
        "picture 48"
    ]

    shapes_to_remove = []

    for shape in slide.shapes:

        name = str(
            getattr(shape, "name", "")
        ).lower()

        if (
            any(name.startswith(p) for p in remove_prefixes)
            or name in remove_names
        ):
            shapes_to_remove.append(
                shape
            )

    for shape in shapes_to_remove:

        element = shape._element

        element.getparent().remove(
            element
        )

        logging.info(
            f"[ppt] Removed "
            f"{shape.name}"
        )




